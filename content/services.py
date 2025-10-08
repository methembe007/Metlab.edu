"""
Content extraction and processing services for the AI learning platform.
"""

import os
import re
import logging
from typing import Optional, Dict, Any
from io import BytesIO

# PDF processing
try:
    import PyPDF2
except ImportError:
    PyPDF2 = None

# OCR capabilities
try:
    from PIL import Image
    import pytesseract
except ImportError:
    Image = None
    pytesseract = None

# Document processing
try:
    from docx import Document
except ImportError:
    Document = None

# PowerPoint processing
try:
    from pptx import Presentation
except ImportError:
    Presentation = None

from django.core.files.base import ContentFile
from .models import UploadedContent

logger = logging.getLogger(__name__)


class ContentExtractionError(Exception):
    """Custom exception for content extraction errors"""
    pass


class ContentExtractor:
    """Main service class for extracting text from various file types"""
    
    def __init__(self):
        self.supported_types = {
            'pdf': self._extract_from_pdf,
            'docx': self._extract_from_docx,
            'txt': self._extract_from_txt,
            'image': self._extract_from_image,
            'pptx': self._extract_from_pptx,
        }
    
    def extract_text(self, uploaded_content: UploadedContent) -> str:
        """
        Extract text from uploaded content based on file type
        
        Args:
            uploaded_content: UploadedContent instance
            
        Returns:
            str: Extracted text content
            
        Raises:
            ContentExtractionError: If extraction fails
        """
        try:
            content_type = uploaded_content.content_type
            
            if content_type not in self.supported_types:
                raise ContentExtractionError(f"Unsupported content type: {content_type}")
            
            extraction_method = self.supported_types[content_type]
            extracted_text = extraction_method(uploaded_content.file)
            
            # Clean and preprocess the extracted text
            cleaned_text = self._preprocess_text(extracted_text)
            
            logger.info(f"Successfully extracted {len(cleaned_text)} characters from {uploaded_content.original_filename}")
            return cleaned_text
            
        except Exception as e:
            logger.error(f"Error extracting text from {uploaded_content.original_filename}: {str(e)}")
            raise ContentExtractionError(f"Failed to extract text: {str(e)}")
    
    def _extract_from_pdf(self, file) -> str:
        """Extract text from PDF file"""
        if PyPDF2 is None:
            raise ContentExtractionError("PyPDF2 not installed. Cannot process PDF files.")
        
        try:
            file.seek(0)
            pdf_reader = PyPDF2.PdfReader(file)
            text_content = []
            
            for page_num, page in enumerate(pdf_reader.pages):
                try:
                    page_text = page.extract_text()
                    if page_text.strip():
                        text_content.append(page_text)
                except Exception as e:
                    logger.warning(f"Error extracting text from page {page_num + 1}: {str(e)}")
                    continue
            
            return '\n\n'.join(text_content)
            
        except Exception as e:
            raise ContentExtractionError(f"PDF extraction failed: {str(e)}")
    
    def _extract_from_docx(self, file) -> str:
        """Extract text from DOCX file"""
        if Document is None:
            raise ContentExtractionError("python-docx not installed. Cannot process DOCX files.")
        
        try:
            file.seek(0)
            doc = Document(file)
            text_content = []
            
            # Extract text from paragraphs
            for paragraph in doc.paragraphs:
                if paragraph.text.strip():
                    text_content.append(paragraph.text)
            
            # Extract text from tables
            for table in doc.tables:
                for row in table.rows:
                    row_text = []
                    for cell in row.cells:
                        if cell.text.strip():
                            row_text.append(cell.text.strip())
                    if row_text:
                        text_content.append(' | '.join(row_text))
            
            return '\n\n'.join(text_content)
            
        except Exception as e:
            raise ContentExtractionError(f"DOCX extraction failed: {str(e)}")
    
    def _extract_from_txt(self, file) -> str:
        """Extract text from TXT file"""
        try:
            file.seek(0)
            content = file.read()
            
            # Try to decode as UTF-8, fallback to other encodings
            if isinstance(content, bytes):
                try:
                    return content.decode('utf-8')
                except UnicodeDecodeError:
                    try:
                        return content.decode('latin-1')
                    except UnicodeDecodeError:
                        return content.decode('utf-8', errors='ignore')
            
            return content
            
        except Exception as e:
            raise ContentExtractionError(f"TXT extraction failed: {str(e)}")
    
    def _extract_from_image(self, file) -> str:
        """Extract text from image using OCR"""
        if Image is None or pytesseract is None:
            raise ContentExtractionError("PIL and pytesseract not installed. Cannot process image files.")
        
        try:
            file.seek(0)
            image = Image.open(file)
            
            # Convert to RGB if necessary
            if image.mode != 'RGB':
                image = image.convert('RGB')
            
            # Extract text using OCR
            extracted_text = pytesseract.image_to_string(image)
            
            return extracted_text
            
        except Exception as e:
            raise ContentExtractionError(f"Image OCR extraction failed: {str(e)}")
    
    def _extract_from_pptx(self, file) -> str:
        """Extract text from PowerPoint file"""
        if Presentation is None:
            raise ContentExtractionError("python-pptx not installed. Cannot process PPTX files.")
        
        try:
            file.seek(0)
            prs = Presentation(file)
            text_content = []
            
            for slide_num, slide in enumerate(prs.slides):
                slide_text = []
                
                # Extract text from shapes
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_text.append(shape.text.strip())
                
                if slide_text:
                    text_content.append(f"Slide {slide_num + 1}:\n" + '\n'.join(slide_text))
            
            return '\n\n'.join(text_content)
            
        except Exception as e:
            raise ContentExtractionError(f"PPTX extraction failed: {str(e)}")
    
    def _preprocess_text(self, text: str) -> str:
        """
        Clean and preprocess extracted text
        
        Args:
            text: Raw extracted text
            
        Returns:
            str: Cleaned and preprocessed text
        """
        if not text:
            return ""
        
        # Remove excessive whitespace
        text = re.sub(r'\s+', ' ', text)
        
        # Remove special characters but keep punctuation
        text = re.sub(r'[^\w\s\.\,\!\?\;\:\-\(\)\[\]\{\}\"\'\/]', '', text)
        
        # Remove multiple consecutive punctuation marks
        text = re.sub(r'([.!?]){2,}', r'\1', text)
        
        # Fix spacing around punctuation
        text = re.sub(r'\s+([.!?,:;])', r'\1', text)
        text = re.sub(r'([.!?])\s*([A-Z])', r'\1 \2', text)
        
        # Remove lines with only numbers or special characters
        lines = text.split('\n')
        cleaned_lines = []
        for line in lines:
            line = line.strip()
            if len(line) > 3 and not re.match(r'^[\d\s\-\.]+$', line):
                cleaned_lines.append(line)
        
        text = '\n'.join(cleaned_lines)
        
        # Final cleanup
        text = text.strip()
        
        return text


class ContentProcessor:
    """Service for processing and analyzing extracted content"""
    
    def __init__(self):
        self.extractor = ContentExtractor()
    
    def process_content(self, uploaded_content: UploadedContent) -> Dict[str, Any]:
        """
        Process uploaded content and extract metadata
        
        Args:
            uploaded_content: UploadedContent instance
            
        Returns:
            dict: Processing results with extracted text and metadata
        """
        try:
            # Extract text content
            extracted_text = self.extractor.extract_text(uploaded_content)
            
            # Analyze content
            analysis = self._analyze_content(extracted_text)
            
            # Update the uploaded content record
            uploaded_content.extracted_text = extracted_text
            uploaded_content.processing_status = 'completed'
            uploaded_content.save()
            
            result = {
                'success': True,
                'extracted_text': extracted_text,
                'word_count': analysis['word_count'],
                'estimated_reading_time': analysis['estimated_reading_time'],
                'language': analysis.get('language', 'en'),
                'complexity_score': analysis.get('complexity_score', 0.5),
            }
            
            logger.info(f"Successfully processed content: {uploaded_content.original_filename}")
            return result
            
        except Exception as e:
            # Mark as failed
            uploaded_content.processing_status = 'failed'
            uploaded_content.save()
            
            logger.error(f"Content processing failed for {uploaded_content.original_filename}: {str(e)}")
            return {
                'success': False,
                'error': str(e),
                'extracted_text': '',
            }
    
    def _analyze_content(self, text: str) -> Dict[str, Any]:
        """
        Analyze extracted text content
        
        Args:
            text: Extracted text content
            
        Returns:
            dict: Analysis results
        """
        if not text:
            return {
                'word_count': 0,
                'estimated_reading_time': 0,
                'complexity_score': 0.0,
            }
        
        # Basic text analysis
        words = text.split()
        word_count = len(words)
        
        # Estimate reading time (average 200 words per minute)
        estimated_reading_time = max(1, word_count // 200)
        
        # Simple complexity score based on average word length and sentence length
        avg_word_length = sum(len(word) for word in words) / max(1, word_count)
        sentences = re.split(r'[.!?]+', text)
        avg_sentence_length = word_count / max(1, len([s for s in sentences if s.strip()]))
        
        # Normalize complexity score (0.0 to 1.0)
        complexity_score = min(1.0, (avg_word_length * 0.1 + avg_sentence_length * 0.05) / 10)
        
        return {
            'word_count': word_count,
            'estimated_reading_time': estimated_reading_time,
            'complexity_score': complexity_score,
            'avg_word_length': avg_word_length,
            'avg_sentence_length': avg_sentence_length,
        }


# Global instances
content_extractor = ContentExtractor()
content_processor = ContentProcessor()